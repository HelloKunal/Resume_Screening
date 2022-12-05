from pptx import Presentation
import aspose.slides as slides
import aspose.pydrawing as drawing


# Instantiate a Presentation object that represents a PPT file

# Save the presentation as PPTX
def process_image(iamge_name):
    if iamge_name.lower().endswith(('ppt')):
        presentation = slides.Presentation(iamge_name)
        iamge_name = iamge_name + 'x'
        presentation.save(iamge_name, slides.export.SaveFormat.PPTX)
    prs = Presentation(iamge_name)
    full_text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text = full_text + shape.text
    return full_text

def print_data(data):
	print(data)

# def output_file(filename, data):
# 	file = open(filename, "w+")
# 	file.write(data)
# 	file.close()

def ppt_to_string(document_path):
	data = process_image(document_path)
	print_data(data)
	return data
